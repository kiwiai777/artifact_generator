#!/usr/bin/env python3
"""
Artifact Generator PoC v3 — metadata output mode.
Solves Dify UTF-8 binary corruption by returning JSON metadata instead of raw binary.
    POST /generate/{fmt}          → binary (legacy, for direct HTTP testing)
    POST /generate/{fmt}/metadata → JSON artifact contract (Dify-compatible)
    GET  /download/{fmt}/{fname}  → serve stored file
    GET  /health                  → status
"""
import json, io, os, sys, argparse, base64, tempfile, hashlib, shutil
from http.server import HTTPServer, BaseHTTPRequestHandler
from urllib.parse import unquote

POC = "/opt/kiwiai/poc/artifact_generator_poc"
OUTPUT_DIR = f"{POC}/output"
TEMPLATES = {
    "pptx": f"{POC}/templates/synthetic_template.pptx",
    "docx": f"{POC}/templates/synthetic_template.docx",
    "xlsx": f"{POC}/templates/synthetic_template.xlsx",
}
ALLOWED_TEMPLATE_DIRS = [f"{POC}/templates", tempfile.gettempdir()]
os.makedirs(OUTPUT_DIR, exist_ok=True)

MIME_MAP = {
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
}

# ---------- Template Mode Resolver ----------
def resolve_template(body: dict, fmt: str):
    mode = body.get("template_mode", "default")
    if mode == "user_upload":
        b64 = body.get("template_base64", "")
        if not b64: raise ValueError("user_upload requires template_base64")
        raw = base64.b64decode(b64)
        suffix = ".pptx" if fmt == "pptx" else ".docx" if fmt == "docx" else ".xlsx"
        tmp = tempfile.NamedTemporaryFile(suffix=suffix, delete=False)
        tmp.write(raw); tmp.close()
        return tmp.name, "[User Upload Template]"
    elif mode == "default":
        tp = body.get("template_path") or TEMPLATES.get(fmt)
        if not tp or not os.path.exists(tp):
            raise ValueError(f"default template not found: {tp}")
        if not any(os.path.realpath(tp).startswith(os.path.realpath(d)) for d in ALLOWED_TEMPLATE_DIRS):
            raise ValueError(f"template path outside allowed dirs: {tp}")
        return tp, "[Default Template]"
    elif mode == "no_template":
        return None, "[No Template]"
    else:
        raise ValueError(f"unknown template_mode: {mode}")

# ---------- Generators ----------
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def generate_pptx(content, template_path, marker):
    prs = Presentation(template_path) if template_path else Presentation()
    sections = content.get("sections", [{"heading": content.get("title","Untitled"), "body": "", "bullets": []}])
    for i, sec in enumerate(sections):
        if i == 0:
            if template_path:
                sl = prs.slides.add_slide(prs.slide_layouts[1] if len(prs.slide_layouts)>1 else prs.slide_layouts[0])
                try: sl.shapes.title.text = marker + " — " + sec["heading"]
                except: pass
            else:
                sl = prs.slides.add_slide(prs.slide_layouts[6])
                tb = sl.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(3))
                tf = tb.text_frame
                p = tf.add_paragraph(); p.text = marker; p.font.size = Pt(20)
                p.font.color.rgb = RGBColor(0x99,0x99,0x99); p.alignment = PP_ALIGN.CENTER
                p = tf.add_paragraph(); p.text = sec["heading"]; p.font.size = Pt(36); p.alignment = PP_ALIGN.CENTER
        else:
            sl = prs.slides.add_slide(prs.slide_layouts[1])
            try: sl.shapes.title.text = sec["heading"]
            except: pass
            try:
                bf = sl.placeholders[1].text_frame
                if sec.get("body"):
                    p = bf.add_paragraph(); p.text = sec["body"]; p.font.size = Pt(18)
                for b in sec.get("bullets", []):
                    p = bf.add_paragraph(); p.text = "• " + b; p.font.size = Pt(16); p.space_after = Pt(8)
            except: pass
    for sl in prs.slides:
        try:
            tb = sl.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12), Inches(0.4))
            tf = tb.text_frame; p = tf.add_paragraph()
            p.text = marker; p.font.size = Pt(8); p.font.color.rgb = RGBColor(0x99,0x99,0x99)
        except: pass
    buf = io.BytesIO(); prs.save(buf); return buf.getvalue()

from docx import Document
from docx.shared import Pt as DPt
from docx.enum.text import WD_ALIGN_PARAGRAPH

def generate_docx(content, template_path, marker):
    doc = Document(template_path) if template_path else Document()
    marker_para = doc.add_paragraph(marker)
    marker_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sections = content.get("sections", [{"heading": content.get("title","Untitled"), "body": "", "bullets": []}])
    for sec in sections:
        doc.add_heading(sec["heading"], level=1)
        if sec.get("body"): doc.add_paragraph(sec["body"])
        for b in sec.get("bullets", []): doc.add_paragraph(b, style="List Bullet")
    buf = io.BytesIO(); doc.save(buf); return buf.getvalue()

from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill

def generate_xlsx(content, template_path, marker):
    wb = Workbook() if not template_path else __import__("openpyxl").load_workbook(template_path)
    ws = wb.active
    ws.title = (content.get("title","Data")[:28] + "-" + content.get("template_mode","default")[:2])[:31]
    ws.insert_rows(1)
    ws.cell(1,1,marker).font = Font(size=10,color="999999")
    if content.get("table_data"):
        td = content["table_data"]
        start = ws.max_row + 1
        hf = Font(bold=True,size=12,color="FFFFFF")
        hf2 = PatternFill(start_color="1A5276",end_color="1A5276",fill_type="solid")
        for i,h in enumerate(td["headers"],1):
            c = ws.cell(row=start,column=i,value=h); c.font=hf; c.fill=hf2
        for ri,row in enumerate(td["rows"]):
            for ci,v in enumerate(row):
                ws.cell(row=start+ri+1,column=ci+1,value=v)
    buf = io.BytesIO(); wb.save(buf); return buf.getvalue()

GENERATORS = {"pptx": generate_pptx, "docx": generate_docx, "xlsx": generate_xlsx}

# ---------- HTTP Server ----------
class Handler(BaseHTTPRequestHandler):
    def log_message(self, fmt, *args): pass

    def do_GET(self):
        path = unquote(self.path)
        if path == "/health":
            return self._json({"status":"ok","service":"artifact-generator-poc-v3"})
        # Download: /download/{fmt}/{filename}
        parts = path.split("/")
        if len(parts) >= 4 and parts[1] == "download":
            fmt, fname = parts[2], parts[3]
            fpath = f"{OUTPUT_DIR}/{fname}"
            if not os.path.realpath(fpath).startswith(os.path.realpath(OUTPUT_DIR)):
                return self._json({"error":"path traversal blocked"}, 403)
            if not os.path.exists(fpath):
                return self._json({"error":"file not found"}, 404)
            data = open(fpath, "rb").read()
            self.send_response(200)
            self.send_header("Content-Type", MIME_MAP.get(fmt, "application/octet-stream"))
            self.send_header("Content-Length", str(len(data)))
            self.end_headers()
            self.wfile.write(data)
        else:
            self._json({"error":"not found"}, 404)

    def do_POST(self):
        path = unquote(self.path)
        parts = path.split("/")
        is_metadata = len(parts) >= 4 and parts[-1] == "metadata"
        fmt = parts[-2] if is_metadata else parts[-1]

        if fmt not in GENERATORS:
            return self._json({"error":f"unknown format: {fmt}"}, 400)

        try:
            length = int(self.headers.get("Content-Length", 0))
            body = json.loads(self.rfile.read(length)) if length > 0 else {}
        except Exception:
            body = {}

        try:
            tmpl_path, marker = resolve_template(body, fmt)
            raw = GENERATORS[fmt](body, tmpl_path, marker)
        except ValueError as e:
            return self._json({"error": str(e)}, 400)
        except Exception as e:
            return self._json({"error": str(e)}, 500)

        if is_metadata:
            # Store file, return artifact contract JSON
            sha = hashlib.sha256(raw).hexdigest()
            fname = f"sandbox-{fmt}-{sha[:12]}.{fmt}"
            fpath = f"{OUTPUT_DIR}/{fname}"
            with open(fpath, "wb") as f:
                f.write(raw)

            title = body.get("title", "Untitled")
            contract = {
                "artifact": {
                    "provider": "external_tool",
                    "provider_file_id": f"{sha[:12]}",
                    "download_url": f"http://172.17.0.1:9191/download/{fmt}/{fname}",
                    "filename": f"{title}.{fmt}",
                    "mime_type": MIME_MAP[fmt],
                    "artifact_type": fmt,
                    "size": len(raw),
                    "checksum": f"sha256:{sha}",
                }
            }
            self._json(contract)
        else:
            # Binary response (legacy)
            self.send_response(200)
            self.send_header("Content-Type", MIME_MAP[fmt])
            self.send_header("Content-Length", str(len(raw)))
            self.end_headers()
            self.wfile.write(raw)

    def _json(self, data, code=200):
        self.send_response(code)
        self.send_header("Content-Type", "application/json")
        self.end_headers()
        self.wfile.write(json.dumps(data, ensure_ascii=False).encode())

if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("--port", type=int, default=9090)
    args = ap.parse_args()
    srv = HTTPServer(("0.0.0.0", args.port), Handler)
    print(f"Artifact Generator PoC v3 on 0.0.0.0:{args.port}")
    srv.serve_forever()
