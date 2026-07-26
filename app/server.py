#!/usr/bin/env python3
"""
Artifact Generator — production service (POC v3 productionized).

Multi-format (pptx/docx/xlsx) generator. Returns a Dify-compatible artifact
contract (JSON metadata) so KPRX / Dify workflows consume it unchanged.

Endpoints:
    POST /generate/metadata         -> JSON artifact contract; format from body
                                      artifact_type (static URL, Dify-safe)
    POST /generate/{fmt}/metadata   -> JSON artifact contract (path-based)
    POST /generate/{fmt}            -> binary (legacy, direct HTTP testing)
    GET  /download/{fmt}/{fname}    -> serve stored artifact
    GET  /health                    -> {"status":"ok",...}

Contract (unchanged from POC) — artifact metadata fields:
    provider / provider_file_id / download_url / filename / mime_type /
    artifact_type / size / checksum

Behavior is preserved from the POC; productionization changes are:
  * bind host configurable (default 172.17.0.1, NOT 0.0.0.0) — Dify-reachable
  * port 9191 (matches production)
  * threaded server for concurrency
  * output / template dirs configurable via env
"""
import argparse
import base64
import hashlib
import io
import json
import os
import sys
import tempfile
import threading
from http.server import ThreadingHTTPServer, BaseHTTPRequestHandler
from urllib.parse import unquote

# ---------------------------------------------------------------------------
# Configuration (env-overridable; defaults match production scope)
# ---------------------------------------------------------------------------
BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUTPUT_DIR = os.environ.get("AG_OUTPUT_DIR", os.path.join(BASE_DIR, "output"))
TEMPLATES_DIR = os.environ.get("AG_TEMPLATES_DIR", os.path.join(BASE_DIR, "templates"))
BIND_HOST = os.environ.get("AG_BIND_HOST", "172.17.0.1")
# External base URL used in download_url of the contract. In production the
# Dify container reaches the host via 172.17.0.1:9191; overridable for tests.
PUBLIC_BASE = os.environ.get("AG_PUBLIC_BASE", "http://172.17.0.1:9191")
DEFAULT_PORT = 9191

TEMPLATES = {
    "pptx": os.path.join(TEMPLATES_DIR, "synthetic_template.pptx"),
    "docx": os.path.join(TEMPLATES_DIR, "synthetic_template.docx"),
    "xlsx": os.path.join(TEMPLATES_DIR, "synthetic_template.xlsx"),
}
# default template_mode may only read from these roots (path-traversal guard)
ALLOWED_TEMPLATE_DIRS = [TEMPLATES_DIR, tempfile.gettempdir()]
os.makedirs(OUTPUT_DIR, exist_ok=True)

# Customer brand-assets root (multi-tenant: each tenant lives under
# <BRAND_ASSETS_DIR>/<tenant>/deck-brand/<logo>). Used ONLY by the brand render
# path; no tenant/customer literal is hard-coded here (redline).
BRAND_ASSETS_DIR = os.environ.get("AG_BRAND_DIR", "/opt/kiwiai/customer-assets")

MIME_MAP = {
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
}


# ---------------------------------------------------------------------------
# Tolerant body-string normalization
# ---------------------------------------------------------------------------
def _normalize(v):
    """Tolerant normalization for body string values.

    Dify JSON-body pills sometimes wrap a value in a literal pair of braces
    (e.g. ``{docx}``) or pad it with whitespace. Strip leading/trailing
    whitespace, then strip at most ONE matching outer ``{ ... }`` pair, then
    strip again. Only parsing is affected; the contract structure is unchanged.

      _normalize("{docx}")   -> "docx"
      _normalize("  docx ")  -> "docx"
      _normalize("docx")     -> "docx"
      _normalize("{{x}}")    -> "{x}"   (one pair only; still rejected downstream)
      _normalize("{}")       -> ""      (-> missing artifact_type)
    """
    s = str(v).strip()
    if len(s) >= 2 and s.startswith("{") and s.endswith("}"):
        s = s[1:-1].strip()
    return s


# ---------------------------------------------------------------------------
# Template Mode Resolver (behavior preserved from POC)
# ---------------------------------------------------------------------------
def resolve_template(body: dict, fmt: str):
    mode = body.get("template_mode", "default")
    if mode == "user_upload":
        b64 = body.get("template_base64", "")
        if not b64:
            raise ValueError("user_upload requires template_base64")
        raw = base64.b64decode(b64)
        suffix = ".pptx" if fmt == "pptx" else ".docx" if fmt == "docx" else ".xlsx"
        tmp = tempfile.NamedTemporaryFile(suffix=suffix, delete=False)
        tmp.write(raw)
        tmp.close()
        return tmp.name, "[User Upload Template]"
    elif mode == "default":
        tp = body.get("template_path") or TEMPLATES.get(fmt)
        if not tp or not os.path.exists(tp):
            raise ValueError(f"default template not found: {tp}")
        if not any(
            os.path.realpath(tp).startswith(os.path.realpath(d) + os.sep)
            or os.path.realpath(tp) == os.path.realpath(d)
            for d in ALLOWED_TEMPLATE_DIRS
        ):
            raise ValueError(f"template path outside allowed dirs: {tp}")
        return tp, "[Default Template]"
    elif mode == "no_template":
        return None, "[No Template]"
    else:
        raise ValueError(f"unknown template_mode: {mode}")


# ---------------------------------------------------------------------------
# Generators (logic preserved from POC)
# ---------------------------------------------------------------------------
# ---------------------------------------------------------------------------
# Brand renderers (IR-adaptive): pptx/docx from a deck IR (brand_tokens + slides).
# Import is tolerant so the server runs both as a package (tests) and as a direct
# script (systemd ExecStart: python app/server.py).
# ---------------------------------------------------------------------------
try:
    from .brand_render import render_pptx_from_ir, render_docx_from_ir
except ImportError:  # direct script run
    import sys as _sys, os as _os
    _sys.path.insert(0, _os.path.dirname(_os.path.abspath(__file__)))
    from brand_render import render_pptx_from_ir, render_docx_from_ir


def _is_brand_ir(content):
    """A body is a deck IR (brand path) iff it has brand_tokens + a slides[] list.
    Everything else (legacy {sections[]} / {table_data}) keeps the synthetic path."""
    return (
        isinstance(content, dict)
        and bool(content.get("brand_tokens"))
        and isinstance(content.get("slides"), list)
    )


def _resolve_brand_dir(content):
    """Resolve the customer deck-brand dir (where logo.png lives) from the body.
    Multi-tenant, no hard-coded slug: body.brand_dir >
    <BRAND_ASSETS_DIR>/<tenant>/deck-brand. Returns None when nothing usable is
    supplied (brand render then proceeds with no logo)."""
    explicit = content.get("brand_dir")
    if explicit:
        return explicit
    tenant = content.get("tenant")
    if tenant and BRAND_ASSETS_DIR:
        return os.path.join(BRAND_ASSETS_DIR, str(tenant), "deck-brand")
    return None


from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def generate_pptx(content, template_path, marker):
    # IR-adaptive: a deck IR (brand_tokens + slides) -> branded pptx; else legacy.
    if _is_brand_ir(content):
        return render_pptx_from_ir(
            content,
            brand_dir=_resolve_brand_dir(content),
            logo_path=content.get("logo_path"),
        )
    prs = Presentation(template_path) if template_path else Presentation()
    sections = content.get(
        "sections", [{"heading": content.get("title", "Untitled"), "body": "", "bullets": []}]
    )
    for i, sec in enumerate(sections):
        if i == 0:
            if template_path:
                sl = prs.slides.add_slide(
                    prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
                )
                try:
                    sl.shapes.title.text = marker + " — " + sec["heading"]
                except Exception:
                    pass
            else:
                sl = prs.slides.add_slide(prs.slide_layouts[6])
                tb = sl.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(3))
                tf = tb.text_frame
                p = tf.add_paragraph()
                p.text = marker
                p.font.size = Pt(20)
                p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
                p.alignment = PP_ALIGN.CENTER
                p = tf.add_paragraph()
                p.text = sec["heading"]
                p.font.size = Pt(36)
                p.alignment = PP_ALIGN.CENTER
        else:
            sl = prs.slides.add_slide(prs.slide_layouts[1])
            try:
                sl.shapes.title.text = sec["heading"]
            except Exception:
                pass
            try:
                bf = sl.placeholders[1].text_frame
                if sec.get("body"):
                    p = bf.add_paragraph()
                    p.text = sec["body"]
                    p.font.size = Pt(18)
                for b in sec.get("bullets", []):
                    p = bf.add_paragraph()
                    p.text = "• " + b
                    p.font.size = Pt(16)
                    p.space_after = Pt(8)
            except Exception:
                pass
    for sl in prs.slides:
        try:
            tb = sl.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12), Inches(0.4))
            tf = tb.text_frame
            p = tf.add_paragraph()
            p.text = marker
            p.font.size = Pt(8)
            p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
        except Exception:
            pass
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


from docx import Document
from docx.shared import Pt as DPt  # noqa: F401  (preserved import for parity)
from docx.enum.text import WD_ALIGN_PARAGRAPH


def generate_docx(content, template_path, marker):
    # IR-adaptive: a deck IR (brand_tokens + slides) -> branded docx; else legacy.
    if _is_brand_ir(content):
        return render_docx_from_ir(
            content,
            brand_dir=_resolve_brand_dir(content),
            work_dir=content.get("work_dir"),
        )
    doc = Document(template_path) if template_path else Document()
    marker_para = doc.add_paragraph(marker)
    marker_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sections = content.get(
        "sections", [{"heading": content.get("title", "Untitled"), "body": "", "bullets": []}]
    )
    for sec in sections:
        doc.add_heading(sec["heading"], level=1)
        if sec.get("body"):
            doc.add_paragraph(sec["body"])
        for b in sec.get("bullets", []):
            doc.add_paragraph(b, style="List Bullet")
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


from openpyxl import Workbook, load_workbook
from openpyxl.styles import Font, PatternFill


def generate_xlsx(content, template_path, marker):
    wb = load_workbook(template_path) if template_path else Workbook()
    ws = wb.active
    ws.title = (content.get("title", "Data")[:28] + "-" + content.get("template_mode", "default")[:2])[:31]
    ws.insert_rows(1)
    ws.cell(1, 1, marker).font = Font(size=10, color="999999")
    if content.get("table_data"):
        td = content["table_data"]
        start = ws.max_row + 1
        hf = Font(bold=True, size=12, color="FFFFFF")
        hf2 = PatternFill(start_color="1A5276", end_color="1A5276", fill_type="solid")
        for i, h in enumerate(td["headers"], 1):
            c = ws.cell(row=start, column=i, value=h)
            c.font = hf
            c.fill = hf2
        for ri, row in enumerate(td["rows"]):
            for ci, v in enumerate(row):
                ws.cell(row=start + ri + 1, column=ci + 1, value=v)
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


GENERATORS = {"pptx": generate_pptx, "docx": generate_docx, "xlsx": generate_xlsx}

# OOXML formats are ZIP containers -> magic bytes "PK\x03\x04"
MAGIC = b"PK\x03\x04"


# ---------------------------------------------------------------------------
# HTTP Server
# ---------------------------------------------------------------------------
class Handler(BaseHTTPRequestHandler):
    server_version = "ArtifactGenerator/1.0"

    def log_message(self, fmt, *args):
        sys.stderr.write("%s - %s\n" % (self.address_string(), fmt % args))

    # ---- helpers ----
    def _json(self, data, code=200):
        body = json.dumps(data, ensure_ascii=False).encode()
        self.send_response(code)
        self.send_header("Content-Type", "application/json; charset=utf-8")
        self.send_header("Content-Length", str(len(body)))
        self.end_headers()
        self.wfile.write(body)

    def _read_json_body(self):
        length = int(self.headers.get("Content-Length", 0))
        if length <= 0:
            return {}
        raw = self.rfile.read(length)
        return json.loads(raw)

    # ---- GET ----
    def do_GET(self):
        path = unquote(self.path)
        if path == "/health":
            return self._json(
                {"status": "ok", "service": "artifact-generator", "version": "1.0"}
            )
        parts = [p for p in path.split("/") if p != ""]
        # /download/{fmt}/{filename}
        if len(parts) >= 3 and parts[0] == "download":
            fmt, fname = parts[1], "/".join(parts[2:])
            fpath = os.path.join(OUTPUT_DIR, fname)
            # path-traversal guard: resolved path must stay under OUTPUT_DIR
            if not (
                os.path.realpath(fpath) == os.path.realpath(OUTPUT_DIR)
                or os.path.realpath(fpath).startswith(os.path.realpath(OUTPUT_DIR) + os.sep)
            ):
                return self._json({"error": "path traversal blocked"}, 403)
            if not os.path.exists(fpath):
                return self._json({"error": "file not found"}, 404)
            with open(fpath, "rb") as f:
                data = f.read()
            self.send_response(200)
            self.send_header("Content-Type", MIME_MAP.get(fmt, "application/octet-stream"))
            self.send_header("Content-Length", str(len(data)))
            self.end_headers()
            self.wfile.write(data)
            return
        return self._json({"error": "not found"}, 404)

    # ---- shared contract builder (path & static routes are byte-identical) ----
    def _metadata_contract(self, body, fmt, raw, normalize_title=False):
        """Store artifact and build the artifact contract dict.

        Used by BOTH `/generate/{fmt}/metadata` and `/generate/metadata` so the
        returned JSON is identical for the same (body, fmt) — same field order,
        same values, same sha-derived ids. When ``normalize_title`` is True
        (static endpoint), the title is run through :func:`_normalize` so Dify's
        brace-wrapped values don't produce dirty filenames like ``{docx}.docx``.
        """
        sha = hashlib.sha256(raw).hexdigest()
        fname = f"sandbox-{fmt}-{sha[:12]}.{fmt}"
        fpath = os.path.join(OUTPUT_DIR, fname)
        with open(fpath, "wb") as f:
            f.write(raw)
        title = body.get("title", "Untitled")
        if normalize_title:
            title = _normalize(title) or "Untitled"
        return {
            "artifact": {
                "provider": "external_tool",
                "provider_file_id": sha[:12],
                "download_url": f"{PUBLIC_BASE}/download/{fmt}/{fname}",
                "filename": f"{title}.{fmt}",
                "mime_type": MIME_MAP[fmt],
                "artifact_type": fmt,
                "size": len(raw),
                "checksum": f"sha256:{sha}",
            }
        }

    def _generate_raw(self, body, fmt):
        """Resolve template + render -> raw bytes. Raises ValueError on bad input."""
        tmpl_path, marker = resolve_template(body, fmt)
        return GENERATORS[fmt](body, tmpl_path, marker)

    # ---- POST /generate/metadata (static URL, format from body) ----
    def _handle_static_metadata(self):
        try:
            body = self._read_json_body()
        except Exception:
            return self._json({"error": "invalid JSON body"}, 400)
        fmt = _normalize(body.get("artifact_type", ""))
        if not fmt:
            return self._json({"error": "missing artifact_type"}, 400)
        if fmt not in GENERATORS:
            return self._json({"error": f"unknown artifact_type: {fmt}"}, 400)
        try:
            raw = self._generate_raw(body, fmt)
        except ValueError as e:
            return self._json({"error": str(e)}, 400)
        except Exception as e:
            return self._json({"error": f"generation failed: {e}"}, 500)
        # identical contract builder as the path-based metadata route
        # (normalize_title=True so Dify brace values don't leak into filename)
        return self._json(self._metadata_contract(body, fmt, raw, normalize_title=True))

    # ---- POST ----
    def do_POST(self):
        path = unquote(self.path)
        parts = [p for p in path.split("/") if p != ""]
        if len(parts) < 2 or parts[0] != "generate":
            return self._json({"error": "not found"}, 404)

        # Static body-format endpoint: POST /generate/metadata
        # (format selected from body.artifact_type; URL has no variable segment,
        #  works around Dify dropping static suffixes after path variables).
        if parts == ["generate", "metadata"]:
            return self._handle_static_metadata()

        # Path-format routes: /generate/{fmt}[/metadata]
        is_metadata = len(parts) >= 3 and parts[-1] == "metadata"
        fmt = parts[-2] if is_metadata else parts[-1]

        if fmt not in GENERATORS:
            return self._json({"error": f"unknown format: {fmt}"}, 400)

        try:
            body = self._read_json_body()
        except Exception:
            return self._json({"error": "invalid JSON body"}, 400)

        try:
            raw = self._generate_raw(body, fmt)
        except ValueError as e:
            return self._json({"error": str(e)}, 400)
        except Exception as e:
            return self._json({"error": f"generation failed: {e}"}, 500)

        if is_metadata:
            return self._json(self._metadata_contract(body, fmt, raw))

        # binary (legacy)
        self.send_response(200)
        self.send_header("Content-Type", MIME_MAP[fmt])
        self.send_header("Content-Length", str(len(raw)))
        self.end_headers()
        self.wfile.write(raw)


def main():
    ap = argparse.ArgumentParser(description="Artifact Generator production service")
    ap.add_argument("--port", type=int, default=DEFAULT_PORT)
    ap.add_argument(
        "--host",
        default=BIND_HOST,
        help="bind host (default 172.17.0.1; Dify-reachable, not 0.0.0.0)",
    )
    args = ap.parse_args()
    srv = ThreadingHTTPServer((args.host, args.port), Handler)
    print(f"Artifact Generator on {args.host}:{args.port} (output={OUTPUT_DIR})", flush=True)
    try:
        srv.serve_forever()
    except KeyboardInterrupt:
        srv.shutdown()


if __name__ == "__main__":
    main()
