#!/usr/bin/env python3
"""
PPTX Native Renderer: IR JSON → .pptx per KIWIAI-DECK-GENERATION-STANDARD-V1 §2.5
Quality targets: text boxes editable, tables native, charts native,
complex diagrams as high-res images.
"""
import json, sys, os, io, tempfile, re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
SLIDE_H_IN = 7.5  # canvas height in inches (for stack vertical-centering math)
FONT_HEADING = 'Microsoft YaHei'
FONT_BODY = 'Microsoft YaHei'

def hex_to_rgb(h):
    h = h.lstrip('#')
    if len(h) == 3:
        h = h[0]*2 + h[1]*2 + h[2]*2
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# SVG helpers (svg_to_png + diagram generators) live in the shared module so both
# the pptx and docx brand renderers use one environment-adaptive render path.
try:
    from .svg_render import svg_to_png, svg_layers, svg_flow
except ImportError:  # direct CLI run (python render_pptx.py ...)
    import sys as _sys
    _sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
    from svg_render import svg_to_png, svg_layers, svg_flow


def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False,
                 color=None, font_name=None, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.name = font_name or FONT_BODY
    p.alignment = alignment
    if color:
        p.font.color.rgb = hex_to_rgb(color)
    return txBox

def add_heading_bar(slide, tag, title, brand_primary):
    pc = hex_to_rgb(brand_primary)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.4), Inches(0.06), Inches(0.35))
    bar.fill.solid()
    bar.fill.fore_color.rgb = pc
    bar.line.fill.background()
    add_text_box(slide, 0.8, 0.38, 1.5, 0.35, tag, font_size=10, bold=True,
                 color=brand_primary, font_name=FONT_HEADING)
    add_text_box(slide, 2.3, 0.35, 10, 0.5, title, font_size=24, bold=True,
                 color='#1A1A2E', font_name=FONT_HEADING)

def add_bullet_list(slide, left, top, width, height, bullets, color='#333'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    from pptx.oxml.ns import qn
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(13)
        p.font.name = FONT_BODY
        p.font.color.rgb = hex_to_rgb(color)
        p.level = 0
        p.space_after = Pt(8)
        pPr = p._p.get_or_add_pPr()
        buChar = pPr.makeelement(qn('a:buChar'), {'char': '\u25cf'})
        for existing in pPr.findall(qn('a:buChar')):
            pPr.remove(existing)
        pPr.append(buChar)
    return txBox

def _render_to_prs(ir, logo_path=None, brand_dir=None, deck_html=None, work_dir=None):
    """Core renderer: IR dict -> pptx Presentation object.

    Args:
        ir: IR dict (deck_ir.json structure: brand_tokens / cover_reference / slides).
        logo_path: Explicit absolute path to the logo image (highest priority).
        brand_dir: Directory containing the customer's deck-brand assets (e.g. logo.png).
                   Combined with IR `cover_reference.logo_path` to resolve the logo.
                   No customer slug/name literal lives in this source (P2-3) — the caller
                   supplies the customer directory.
        deck_html: Optional deck HTML (source of SVG diagrams). When absent, diagram
                   slides are rendered from IR data (svg_layers / svg_flow).
        work_dir: Working dir for temp PNG files (default: tempdir/kiwiai-brand-render).
    """
    work_dir = work_dir or os.path.join(tempfile.gettempdir(), 'kiwiai-brand-render')
    os.makedirs(work_dir, exist_ok=True)

    # Read logo path from IR and resolve to absolute path (per P2-3: no hardcoded slug).
    # Resolution order: explicit logo_path param > brand_dir + IR logo_path > unresolved (no logo).
    logo_rel = ir.get("cover_reference", {}).get("logo_path", "")
    if not logo_path and logo_rel and brand_dir:
        logo_path = os.path.join(brand_dir, logo_rel)
    
    
    brand = ir['brand_tokens']
    primary = brand['palette']['primary']
    primary_dark = brand['palette'].get('primary_dark', '#005A9B')
    primary_light = brand['palette']['primary_light']
    primary_pale = brand['palette']['primary_pale']
    accent = brand['palette']['accent']
    # Extract company short name from full CN name: strip location prefix + legal suffix
    company_cn_full = ir['cover_reference']['company_name_cn']
    import re as _re
    _short = _re.sub(r'^(北京|上海|深圳|广州|杭州|成都|武汉|南京|重庆|天津|苏州|西安|长沙|郑州|青岛|大连|厦门|宁波)', '', company_cn_full)
    _short = _re.sub(r'(科技|技术|网络|信息|数据|软件|智能)?(有限公司|股份公司|集团公司|有限责任公司|股份有限公司|集团有限公司)$', '', _short)
    company_short_name = _short.strip() or company_cn_full
    
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]
    temp_files = []
    
    for idx, slide_data in enumerate(ir['slides']):
        slide = prs.slides.add_slide(blank_layout)
        stype = slide_data['type']
        background = slide.background
        fill = background.fill
        fill.solid()
        
        if stype == 'cover':
            bg_color = brand['cover'].get('bg_color', '#F2F3F5')
            fill.fore_color.rgb = hex_to_rgb(bg_color)
            # Vertically center the cover stack (logo -> company_cn -> company_en ->
            # divider -> title -> subtitle) on the 7.5" canvas, replicating the golden
            # HTML deck `.slide-cover .slide-inner{justify-content:center}`. Each
            # element sits at `cover_y0 + <offset>`; offsets preserve the original
            # golden spacing, cover_y0 centers the stack (midpoint at SLIDE_H_IN/2).
            # Footer stays pinned at the bottom (unchanged). [KIWIAGENCY-DECK-PPTX-COVER-CENTER-0]
            cover_off_logo, cover_off_company_cn, cover_off_company_en = 0.0, 1.1, 1.45
            cover_off_divider, cover_off_title, cover_off_subtitle = 1.85, 2.1, 3.2
            cover_subtitle_h = 0.5
            cover_stack_half = (cover_off_subtitle + cover_subtitle_h - cover_off_logo) / 2  # 1.85
            cover_y0 = SLIDE_H_IN / 2 - cover_stack_half                              # 1.9
            if logo_path and os.path.exists(logo_path):
                logo_left = (13.333 - 1.5) / 2
                slide.shapes.add_picture(logo_path, Inches(logo_left),
                                         Inches(cover_y0 + cover_off_logo), Inches(1.5), Inches(0.9))
            company_cn = ir['cover_reference']['company_name_cn']
            add_text_box(slide, 1.5, cover_y0 + cover_off_company_cn, 10.333, 0.4, company_cn,
                         font_size=15, bold=True, color='#1A1A2E', font_name=FONT_HEADING, alignment=PP_ALIGN.CENTER)
            company_en = ir['cover_reference']['company_name_en']
            add_text_box(slide, 1.5, cover_y0 + cover_off_company_en, 10.333, 0.3, company_en,
                         font_size=11, color='#666', font_name=FONT_BODY, alignment=PP_ALIGN.CENTER)
            div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.416),
                                         Inches(cover_y0 + cover_off_divider), Inches(0.5), Inches(0.03))
            div.fill.solid()
            div.fill.fore_color.rgb = hex_to_rgb(primary)
            div.line.fill.background()
            title = slide_data.get('title', '')
            add_text_box(slide, 1.0, cover_y0 + cover_off_title, 11.333, 1.2, title,
                         font_size=34, bold=True, color=primary, font_name=FONT_HEADING, alignment=PP_ALIGN.CENTER)
            subtitle = slide_data.get('subtitle', '')
            if subtitle:
                add_text_box(slide, 2.0, cover_y0 + cover_off_subtitle, 9.333, cover_subtitle_h, subtitle,
                             font_size=14, color='#666', font_name=FONT_BODY, alignment=PP_ALIGN.CENTER)
            footer_text = company_short_name + ' · 2026'
            add_text_box(slide, 0, 6.9, 13.333, 0.3, footer_text, font_size=10,
                         color='#999', font_name=FONT_BODY, alignment=PP_ALIGN.CENTER)
        
        elif stype == 'end':
            # Centered layout per KIWIAI-DECK-GENERATION-STANDARD-V1 §2, replicating golden end slide
            # No left accent block, no watermark; all elements horizontally centered on white/light bg
            bg_color = brand['cover'].get('bg_color', '#F2F3F5')
            fill.fore_color.rgb = hex_to_rgb(bg_color)
            slide_w_in = 13.333
            # Vertically center the end stack (logo -> company_cn -> company_en ->
            # divider -> title -> subtitle) on the 7.5" canvas, same as cover (golden
            # `.slide-end .slide-inner{justify-content:center}`). Contact info is pinned
            # to the lower part (贴下部), bottom-anchored at the cover-footer baseline so
            # the block sits cleanly under the centered stack with no overlap (mirrors
            # the cover footer's bottom pin). [KIWIAGENCY-DECK-PPTX-COVER-CENTER-0]
            end_off_logo, end_off_company_cn, end_off_company_en = 0.0, 1.1, 1.45
            end_off_divider, end_off_title, end_off_subtitle = 1.85, 2.1, 3.0
            end_subtitle_h = 0.5
            end_stack_half = (end_off_subtitle + end_subtitle_h - end_off_logo) / 2  # 1.75
            end_y0 = SLIDE_H_IN / 2 - end_stack_half                                # 2.0
            # Logo centered at top
            if logo_path and os.path.exists(logo_path):
                logo_w, logo_h = 1.5, 0.9
                slide.shapes.add_picture(logo_path, Inches((slide_w_in - logo_w) / 2),
                                         Inches(end_y0 + end_off_logo), Inches(logo_w), Inches(logo_h))
            # Company name CN (brand color, centered)
            add_text_box(slide, 1.5, end_y0 + end_off_company_cn, 10.333, 0.4, company_cn_full,
                         font_size=15, bold=True, color=primary, font_name=FONT_HEADING, alignment=PP_ALIGN.CENTER)
            # Company name EN (centered)
            company_en = ir['cover_reference']['company_name_en']
            add_text_box(slide, 1.5, end_y0 + end_off_company_en, 10.333, 0.3, company_en,
                         font_size=11, color='#666', font_name=FONT_BODY, alignment=PP_ALIGN.CENTER)
            # Divider
            div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.416),
                                         Inches(end_y0 + end_off_divider), Inches(0.5), Inches(0.03))
            div.fill.solid()
            div.fill.fore_color.rgb = hex_to_rgb(primary)
            div.line.fill.background()
            # Main title (感谢聆听)
            title = slide_data.get('title', '感谢聆听')
            add_text_box(slide, 1.0, end_y0 + end_off_title, 11.333, 1.0, title, font_size=34,
                         bold=True, color=primary, font_name=FONT_HEADING, alignment=PP_ALIGN.CENTER)
            # Subtitle (方案标题)
            subtitle = slide_data.get('subtitle', '')
            if subtitle:
                add_text_box(slide, 2.0, end_y0 + end_off_subtitle, 9.333, end_subtitle_h, subtitle,
                             font_size=16, color='#666', font_name=FONT_HEADING, alignment=PP_ALIGN.CENTER)
            # Contact info centered, pinned to the lower part (贴下部): bottom-anchored at
            # the cover-footer baseline (6.9+0.3=7.2) so the last line aligns with the
            # footer and the block clears the centered stack (no overlap). line_step keeps
            # the P2-4 >=0.15" inter-line gap.
            contact = slide_data.get('contact', '')
            if contact:
                parts = contact.split('|')
                line_h = 0.35
                line_step = 0.35 + 0.15  # 0.5 (P2-4: >=0.15" gap between lines)
                contact_baseline = 7.2   # last line bottom == cover footer bottom (贴下部)
                y = contact_baseline - line_h - (len(parts) - 1) * line_step
                for part in parts:
                    part = part.strip()
                    add_text_box(slide, 3.0, y, 7.333, line_h, part, font_size=12,
                                 color='#999', font_name=FONT_BODY, alignment=PP_ALIGN.CENTER)
                    y += line_step
        
        elif stype == 'content':
            fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            tag = slide_data.get('tag', '')
            title = slide_data.get('title', '')
            add_heading_bar(slide, tag, title, primary)
            
            if 'agenda_items' in slide_data:
                items = slide_data['agenda_items']
                for i, item in enumerate(items):
                    col = i % 2
                    row = i // 2
                    x = 0.8 + col * 6.2
                    y = 1.3 + row * 1.15
                    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.5), Inches(0.5))
                    circ.fill.solid()
                    circ.fill.fore_color.rgb = hex_to_rgb(primary)
                    circ.line.fill.background()
                    tf = circ.text_frame
                    tf.word_wrap = False
                    p = tf.paragraphs[0]
                    p.text = str(i + 1).zfill(2)
                    p.font.size = Pt(14)
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    p.font.name = FONT_HEADING
                    p.alignment = PP_ALIGN.CENTER
                    add_text_box(slide, x + 0.7, y + 0.08, 5.0, 0.4, item,
                                 font_size=13, bold=True, color='#1A1A2E', font_name=FONT_BODY)
            
            elif 'bullets' in slide_data:
                has_two_col = slide_data.get('layout') == 'two-column'
                if has_two_col:
                    add_bullet_list(slide, 0.7, 1.2, 5.5, 4.5, slide_data['bullets'])
                    hl_x, hl_y, hl_w, hl_h = 7.0, 1.2, 5.8, 4.5
                    hl_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(hl_x), Inches(hl_y), Inches(hl_w), Inches(hl_h))
                    hl_box.fill.solid()
                    hl_box.fill.fore_color.rgb = hex_to_rgb(primary_pale)
                    hl_box.line.color.rgb = hex_to_rgb(primary)
                    hl_box.line.width = Pt(2)
                    accent_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(hl_x), Inches(hl_y), Inches(0.05), Inches(hl_h))
                    accent_line.fill.solid()
                    accent_line.fill.fore_color.rgb = hex_to_rgb(primary)
                    accent_line.line.fill.background()
                    add_text_box(slide, hl_x + 0.3, hl_y + 0.2, hl_w - 0.4, 0.35,
                                 company_short_name + '承诺', font_size=14, bold=True,
                                 color=primary, font_name=FONT_HEADING)
                    hl_text = slide_data.get('highlight', '')
                    for stat_item in slide_data.get('stats', []):
                        hl_text = hl_text.replace(' ' + stat_item['value'] + ' ' + stat_item['label'], '')
                    hl_text = hl_text.replace('  ', ' ').strip()
                    add_text_box(slide, hl_x + 0.3, hl_y + 0.7, hl_w - 0.5, 1.8,
                                 hl_text, font_size=12, color='#333', font_name=FONT_BODY)
                    stats = slide_data.get('stats', [])
                    stat_y = hl_y + 2.8
                    for si, stat in enumerate(stats):
                        sx = hl_x + 0.3 + si * 2.8
                        add_text_box(slide, sx, stat_y, 2.5, 0.45, stat['value'],
                                     font_size=22, bold=True, color=primary, font_name=FONT_HEADING)
                        add_text_box(slide, sx, stat_y + 0.45, 2.5, 0.3, stat['label'],
                                     font_size=10, color='#666', font_name=FONT_BODY)
                else:
                    add_bullet_list(slide, 0.7, 1.2, 12.0, 5.5, slide_data['bullets'])
            
            elif 'cards' in slide_data:
                cards = slide_data['cards']
                n = len(cards)
                card_w = min(2.8, 12.5 / n - 0.3)
                start_x = (13.333 - (card_w + 0.3) * n + 0.3) / 2
                for ci, card in enumerate(cards):
                    cx = start_x + ci * (card_w + 0.3)
                    cy = 1.3
                    card_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx), Inches(cy), Inches(card_w), Inches(4.2))
                    card_shape.fill.solid()
                    card_shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    card_shape.line.color.rgb = hex_to_rgb('#E1E5EB')
                    card_shape.line.width = Pt(1)
                    bar_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx), Inches(cy), Inches(card_w), Inches(0.04))
                    bar_shape.fill.solid()
                    bar_shape.fill.fore_color.rgb = hex_to_rgb(primary)
                    bar_shape.line.fill.background()
                    icon = card.get('icon', '')
                    if icon:
                        add_text_box(slide, cx + 0.3, cy + 0.4, card_w - 0.6, 0.5,
                                     icon, font_size=28, color=primary, font_name=FONT_BODY, alignment=PP_ALIGN.CENTER)
                    card_title = card.get('title', '')
                    add_text_box(slide, cx + 0.2, cy + 1.1, card_w - 0.4, 0.5,
                                 card_title, font_size=15, bold=True, color='#1A1A2E', font_name=FONT_HEADING)
                    body = card.get('body', '')
                    if body:
                        add_text_box(slide, cx + 0.2, cy + 1.7, card_w - 0.4, 1.5,
                                     body, font_size=11, color='#666', font_name=FONT_BODY)
                    stat = card.get('stat', '')
                    stat_label = card.get('stat_label', '')
                    if stat:
                        add_text_box(slide, cx + 0.2, cy + 3.0, card_w - 0.4, 0.5,
                                     stat, font_size=28, bold=True, color=primary, font_name=FONT_HEADING)
                    if stat_label:
                        add_text_box(slide, cx + 0.2, cy + 3.5, card_w - 0.4, 0.3,
                                     stat_label, font_size=10, color='#888', font_name=FONT_BODY)
            
            elif 'timeline' in slide_data:
                timeline = slide_data['timeline']
                n = len(timeline)
                line_y = 3.2
                line_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(line_y), Inches(11.333), Inches(0.02))
                line_shape.fill.solid()
                line_shape.fill.fore_color.rgb = hex_to_rgb(primary_light)
                line_shape.line.fill.background()
                step_w = 10.5 / n
                for ti, item in enumerate(timeline):
                    cx = 1.3 + ti * step_w
                    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx + step_w/2 - 0.13), Inches(line_y - 0.13), Inches(0.26), Inches(0.26))
                    dot.fill.solid()
                    dot.fill.fore_color.rgb = hex_to_rgb(primary)
                    dot.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    dot.line.width = Pt(3)
                    phase = item.get('phase', '')
                    add_text_box(slide, cx, line_y - 0.8, step_w, 0.3, 'Phase ' + phase,
                                 font_size=9, bold=True, color=primary, font_name=FONT_HEADING, alignment=PP_ALIGN.CENTER)
                    ttl = item.get('title', '')
                    add_text_box(slide, cx, line_y + 0.4, step_w, 0.4, ttl,
                                 font_size=14, bold=True, color='#1A1A2E', font_name=FONT_HEADING, alignment=PP_ALIGN.CENTER)
                    desc = item.get('description', '')
                    add_text_box(slide, cx, line_y + 0.9, step_w, 1.5, desc,
                                 font_size=10, color='#666', font_name=FONT_BODY, alignment=PP_ALIGN.CENTER)
            
            elif 'table' in slide_data:
                tbl_ir = slide_data['table']
                headers = tbl_ir.get('headers', [])
                rows = tbl_ir.get('rows', [])
                if headers and rows:
                    n_cols = len(headers)
                    n_rows = len(rows) + 1  # +1 for header row
                    # Native pptx table, sized to fit the content area
                    tbl_left, tbl_top = 0.8, 1.2
                    tbl_w = 11.733  # ~13.333 - 2*0.8 margins
                    row_h = 0.4
                    tbl_h = n_rows * row_h
                    tbl_shape = slide.shapes.add_table(n_rows, n_cols, Inches(tbl_left), Inches(tbl_top), Inches(tbl_w), Inches(tbl_h))
                    table = tbl_shape.table
                    # Column width: evenly distributed
                    col_w = int(Inches(tbl_w) / n_cols)
                    for ci in range(n_cols):
                        table.columns[ci].width = col_w
                    # Header row: brand primary bg + white bold text
                    # python-pptx table.cell(row, col)
                    for ci, hdr in enumerate(headers):
                        cell = table.cell(0, ci)
                        # Set cell fill to primary color
                        cell_fill = cell.fill
                        cell_fill.solid()
                        cell_fill.fore_color.rgb = hex_to_rgb(primary)
                        # Set text
                        tf = cell.text_frame
                        tf.word_wrap = True
                        p = tf.paragraphs[0]
                        p.text = hdr
                        p.font.size = Pt(11)
                        p.font.bold = True
                        p.font.name = FONT_HEADING
                        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                        p.alignment = PP_ALIGN.CENTER
                        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                    # Data rows: body font, alternating row shading
                    for ri, row_data in enumerate(rows):
                        for ci, val in enumerate(row_data):
                            if ci >= n_cols:
                                break
                            cell = table.cell(ri + 1, ci)
                            cell_fill = cell.fill
                            cell_fill.solid()
                            if ri % 2 == 0:
                                cell_fill.fore_color.rgb = hex_to_rgb('F8FAFB')
                            else:
                                cell_fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                            tf = cell.text_frame
                            tf.word_wrap = True
                            p = tf.paragraphs[0]
                            p.text = str(val)
                            p.font.size = Pt(10)
                            p.font.name = FONT_BODY
                            p.font.color.rgb = hex_to_rgb('#333333')
                            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        elif stype == 'diagram':
            fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            tag = slide_data.get('tag', '')
            title = slide_data.get('title', '')
            add_heading_bar(slide, tag, title, primary)
            
            # Resolve SVG: prefer a matching SVG from deck.html (if provided), else
            # generate from IR data (svg_layers / svg_flow) so rendering works without
            # a deck.html dependency (server integration has no deck.html in the body).
            matching_svg = None
            if deck_html:
                for svg_m in re.finditer(r'<svg[^>]*>(.*?)</svg>', deck_html, re.DOTALL):
                    svg_start = svg_m.start()
                    before = deck_html[max(0, svg_start-4000):svg_start]
                    if title in before:
                        matching_svg = svg_m.group(0)
                        break
            if not matching_svg:
                svg_kind = slide_data.get('svg_kind', '')
                if svg_kind == 'layers':
                    layers = slide_data.get('layers', [])
                    if layers:
                        matching_svg = svg_layers(title, layers, primary, primary_pale)
                elif svg_kind == 'flow':
                    steps = slide_data.get('steps', [])
                    if steps:
                        matching_svg = svg_flow(title, steps, primary, primary_light)
            
            if matching_svg:
                png_path = os.path.join(work_dir, '_diagram_' + str(idx) + '.png')
                temp_files.append(png_path)
                native_w, native_h = 1280, 600
                svg_to_png(matching_svg, png_path, width=native_w, height=native_h)
                # svg_to_png raises RuntimeError on any failure (no silent degrade, §2
                # 降级禁令 / ADR-017) -- reaching here means the PNG exists and is valid.
                # P2-1: scale by native aspect ratio (1280x600 = 2.133) instead of stretching
                # into the old fixed 12.333x4.8 (2.569) box; center the scaled image in the box.
                box_left, box_top, box_w, box_h = 0.5, 1.2, 12.333, 4.8
                native_ratio = native_w / native_h
                box_ratio = box_w / box_h
                if native_ratio > box_ratio:
                    disp_w = box_w
                    disp_h = disp_w / native_ratio
                else:
                    disp_h = box_h
                    disp_w = disp_h * native_ratio
                disp_left = box_left + (box_w - disp_w) / 2
                disp_top = box_top + (box_h - disp_h) / 2
                slide.shapes.add_picture(png_path, Inches(disp_left), Inches(disp_top), Inches(disp_w), Inches(disp_h))
            else:
                add_text_box(slide, 2, 3, 9, 2,
                             '[\u56fe\u793a: ' + title + ']\n(\u672a\u627e\u5230\u5bf9\u5e94 SVG)',
                             font_size=14, color='#999', alignment=PP_ALIGN.CENTER)
            
            insight = slide_data.get('insight', '')
            if insight:
                sentences = re.split(r'(?<=[\u3002])\s*', insight)
                sentences = [s.strip() for s in sentences if s.strip()]
                y = 6.3
                for sent in sentences[:2]:
                    add_text_box(slide, 0.7, y, 12.0, 0.3, sent, font_size=9, color='#888', font_name=FONT_BODY)
                    y += 0.22
        
        elif stype == 'chart':
            fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            tag = slide_data.get('tag', '')
            title = slide_data.get('title', '')
            add_heading_bar(slide, tag, title, primary)
            
            labels = slide_data.get('labels', [])
            values = slide_data.get('values', [])
            unit = slide_data.get('unit', '%')
            
            if labels and values:
                chart_data = CategoryChartData()
                chart_data.categories = labels
                chart_data.add_series(title, values)
                chart_frame = slide.shapes.add_chart(
                    XL_CHART_TYPE.COLUMN_CLUSTERED,
                    Inches(1.0), Inches(1.2), Inches(11.333), Inches(4.5),
                    chart_data
                )
                chart = chart_frame.chart
                chart.has_legend = False
                plot = chart.plots[0]
                series = plot.series[0]
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = hex_to_rgb(primary)
                value_axis = chart.value_axis
                value_axis.has_title = True
                value_axis.axis_title.text_frame.paragraphs[0].text = unit
                value_axis.major_gridlines.format.line.color.rgb = hex_to_rgb('#E8E8E8')
                cat_axis = chart.category_axis
                cat_axis.tick_labels.font.size = Pt(9)
                cat_axis.tick_labels.font.name = FONT_BODY
            else:
                add_text_box(slide, 2, 3, 9, 2,
                             '[\u56fe\u8868: ' + title + ']\n(\u65e0\u6570\u636e)',
                             font_size=14, color='#999', alignment=PP_ALIGN.CENTER)
            
            insight = slide_data.get('insight', '')
            if insight:
                sentences = re.split(r'(?<=[\u3002])\s*', insight)
                sentences = [s.strip() for s in sentences if s.strip()]
                y = 6.1
                for sent in sentences[:2]:
                    add_text_box(slide, 0.7, y, 12.0, 0.3, sent, font_size=9, color='#888', font_name=FONT_BODY)
                    y += 0.22
    
    for tf in temp_files:
        try:
            if os.path.exists(tf):
                os.unlink(tf)
        except:
            pass
    return prs


def render_pptx_from_ir(ir, brand_dir=None, logo_path=None, deck_html=None, work_dir=None):
    """Library entry: IR dict -> pptx bytes (used by artifact_generator server)."""
    prs = _render_to_prs(ir, logo_path=logo_path, brand_dir=brand_dir,
                         deck_html=deck_html, work_dir=work_dir)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def render_pptx(ir_path, deck_html_path=None, output_path=None, logo_path=None,
                brand_dir=None, work_dir=None):
    """CLI entry: IR file -> output .pptx file."""
    with open(ir_path, 'r', encoding='utf-8') as f:
        ir = json.load(f)
    deck_html = None
    if deck_html_path and os.path.exists(deck_html_path):
        with open(deck_html_path, 'r', encoding='utf-8') as f:
            deck_html = f.read()
    prs = _render_to_prs(ir, logo_path=logo_path, brand_dir=brand_dir,
                         deck_html=deck_html, work_dir=work_dir)
    output_path = output_path or 'output.pptx'
    prs.save(output_path)
    return prs


if __name__ == '__main__':
    if len(sys.argv) < 2:
        print("Usage: python3 render_pptx.py <deck_ir.json> [deck.html] [output.pptx] [logo_path] [brand_dir]")
        sys.exit(1)
    ir_path = sys.argv[1]
    deck_html = sys.argv[2] if len(sys.argv) > 2 else None
    output = sys.argv[3] if len(sys.argv) > 3 else 'output.pptx'
    logo = sys.argv[4] if len(sys.argv) > 4 else None
    brand_dir = sys.argv[5] if len(sys.argv) > 5 else None
    render_pptx(ir_path, deck_html, output, logo, brand_dir)
    print('PPTX saved: ' + output)
